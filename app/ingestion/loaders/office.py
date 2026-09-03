import os
import logfire


def _parse_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def parse_office(file_path: str) -> str:
    """
    Parses Office documents (.docx, .pptx) directly using lightweight python-docx and python-pptx.
    Falls back to Unstructured if available for legacy compatibility.
    """
    with logfire.span("📄 Office Document Parsing", filename=file_path):
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".docx":
                full_text = _parse_docx(file_path)
            elif ext == ".pptx":
                full_text = _parse_pptx(file_path)
            else:
                try:
                    from unstructured.partition.auto import partition
                    elements = partition(filename=file_path)
                    full_text = "\n".join([str(el) for el in elements])
                except ImportError:
                    logfire.warning(f"Unsupported office format without unstructured: {ext}")
                    return ""

            if not full_text.strip():
                logfire.warning(f"⚠️ Empty text extracted from {file_path}")
            else:
                logfire.info(f"✅ Successfully parsed {len(full_text)} characters from {file_path}")

            return full_text
        except Exception as e:
            logfire.error(f"❌ Office Parse Failed: {e}")
            raise e
